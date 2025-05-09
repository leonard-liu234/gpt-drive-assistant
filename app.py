
from flask import Flask, request, jsonify, send_file, redirect
from flask_cors import CORS
from pptx import Presentation
import os
import requests
from datetime import datetime

app = Flask(__name__, static_folder="static", static_url_path="")
CORS(app)

ACCESS_TOKEN = os.getenv("GDRIVE_ACCESS_TOKEN")

# OAuth proxy endpoints
GOOGLE_AUTH_BASE = "https://accounts.google.com/o/oauth2/v2/auth"
GOOGLE_TOKEN_BASE = "https://oauth2.googleapis.com/token"

@app.route("/oauth/authorize")
def proxy_authorize():
    return redirect(f"{GOOGLE_AUTH_BASE}?{request.query_string.decode()}")

@app.route("/oauth/token", methods=["POST"])
def proxy_token():
    headers = {"Content-Type": "application/x-www-form-urlencoded"}
    response = requests.post(GOOGLE_TOKEN_BASE, data=request.form, headers=headers)
    return response.content, response.status_code, response.headers.items()

@app.route("/", methods=["GET"])
def index():
    return "✅ GPT Drive Assistant is running. Use /generate-ppt or /folders/<folder_id>/list to start."

@app.route("/.well-known/ai-plugin.json", methods=["GET"])
def plugin_manifest():
    return send_file("static/.well-known/ai-plugin.json", mimetype="application/json")

@app.route("/generate-ppt", methods=["POST"])
def generate_ppt():
    data = request.get_json()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "确认函摘要"
    body.text = f"""项目名称：{data.get('project_name')}
客户名称：{data.get('client_name')}
联系方式：{data.get('contact')}
报价编号：{data.get('quote_number')}
报价日期：{data.get('quote_date')}"""
    os.makedirs("generated_ppt", exist_ok=True)
    filename = f"confirmation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    filepath = os.path.join("generated_ppt", filename)
    prs.save(filepath)
    return jsonify({"download_url": f"/download-ppt/{filename}"})

@app.route("/download-ppt/<filename>", methods=["GET"])
def download_ppt(filename):
    return send_file(os.path.join("generated_ppt", filename), as_attachment=True)

@app.route("/folders/<folder_id>/list", methods=["GET"])
def list_folder_files(folder_id):
    def get_files(folder_id):
        headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}
        params = {
            "q": f"'{folder_id}' in parents",
            "fields": "files(id,name,mimeType)",
            "pageSize": 1000
        }
        response = requests.get("https://www.googleapis.com/drive/v3/files", headers=headers, params=params)
        if response.status_code != 200:
            return []
        files = response.json().get("files", [])
        all_files = []
        for f in files:
            all_files.append(f)
            if f.get("mimeType") == "application/vnd.google-apps.folder":
                all_files.extend(get_files(f["id"]))
        return all_files

    return jsonify(get_files(folder_id))

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000)
